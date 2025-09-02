from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Helper to add background image
def add_background_slide(prs, layout, img_path, title_text, body_text=""):
    slide = prs.slides.add_slide(layout)

    # Add picture covering whole slide
    left = top = Inches(0)
    pic = slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)
    # Send picture to back
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    # Body text
    if body_text:
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), prs.slide_width - Inches(1), prs.slide_height - Inches(3))
        tf = body_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = body_text
        run.font.size = Pt(22)
        run.font.color.rgb = RGBColor(240, 240, 240)

    return slide


# === Presentation Builder ===
prs = Presentation()

# Title Slide
add_background_slide(prs, prs.slide_layouts[6], "images/title.jpg",
                     "Global Technological Disruption:\nAI & Automation in Nigeria",
                     "Opportunities and Challenges")

# Content slides
slides_content = [
    ("Introduction", "AI and automation are transforming economies worldwide.\nNigeria is rapidly integrating these technologies."),
    ("Key Drivers", "• AI adoption in banking, agriculture, health\n• Rise of fintech and automation startups\n• Government push for digital economy"),
    ("Challenges", "• Job displacement fears\n• Digital skills gap\n• Infrastructure limitations\n• Cybersecurity risks"),
    ("Economic Impact", "• Growth of fintech sector\n• Boost in digital services GDP contribution\n• Efficiency in logistics, telecom, and health"),
    ("Industry Impact", "Banking: automated services\nAgriculture: AI for crop yield prediction\nHealth: telemedicine & diagnostics\nManufacturing: robotics adoption"),
    ("Opportunities", "• Youth-driven tech ecosystem\n• Startups building African AI\n• Government investment in digital literacy\n• Export of digital talent"),
    ("Way Forward", "• Upskilling workforce\n• Policies for ethical AI\n• Encourage local R&D\n• Infrastructure modernization"),
    ("Conclusion", "AI and automation present both risks and opportunities.\nNigeria must balance innovation with inclusivity.")
]

for title, body in slides_content:
    add_background_slide(prs, prs.slide_layouts[6], "images/bg.jpg", title, body)

# Save
prs.save("AI_Automation_Nigeria_Modern.pptx")
print("✅ Modern presentation saved as AI_Automation_Nigeria_Modern.pptx")
