from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def add_text_slide(prs, title, content, icon_path=None, bg_color=RGBColor(245, 245, 245)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Optional icon
    if icon_path:
        slide.shapes.add_picture(icon_path, Inches(0.3), Inches(0.3), height=Inches(0.7))

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), prs.slide_width - Inches(1), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(40, 40, 40)

    # Body
    body_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), prs.slide_width - Inches(2), prs.slide_height - Inches(2))
    tf = body_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = content
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(60, 60, 60)

    return slide


# === Presentation Builder ===
prs = Presentation()

# Title Slide with background photo
slide = prs.slides.add_slide(prs.slide_layouts[6])
pic = slide.shapes.add_picture("images/title.jpg", Inches(0), Inches(0),
                               width=prs.slide_width, height=prs.slide_height)
slide.shapes._spTree.remove(pic._element)
slide.shapes._spTree.insert(2, pic._element)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), prs.slide_width - Inches(1), Inches(2))
tf = title_box.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "AI & Automation in Nigeria\nDisruption and Opportunities"
run.font.size = Pt(44)
run.font.bold = True
run.font.color.rgb = RGBColor(255, 255, 255)

# Content slides
slides_content = [
    ("Introduction", "AI and automation are reshaping work and industry in Nigeria."),
    ("Drivers of Change", "Fintech boom, digital adoption, and global AI trends."),
    ("Challenges", "Job losses, skills gap, infrastructure, policy uncertainty."),
    ("Economic Impact", "AI boosts fintech, healthcare, logistics, and agriculture."),
    ("Industry Impact", "• Banking automation\n• Telemedicine\n• Smart farming\n• Robotics in factories"),
    ("Opportunities", "Nigeria’s youth, startups, and digital talent pool."),
    ("Way Forward", "Upskilling, regulation, and investment in R&D."),
    ("Conclusion", "Balance disruption with inclusion for sustainable growth.")
]

for i, (title, content) in enumerate(slides_content):
    add_text_slide(prs, title, content, icon_path="icons/icon{}.png".format(i+1))

# Save
prs.save("AI_Automation_Nigeria_Business.pptx")
print("✅ Business presentation saved as AI_Automation_Nigeria_Business.pptx")
